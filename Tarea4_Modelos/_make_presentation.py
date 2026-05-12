# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH

BASE = r"c:\Users\PC\Desktop\Semestre 9\Analítica\Proyecto2\Proyecto2ACTM\Tarea4_Modelos"
PPTX_PATH = os.path.join(BASE, "presentacion_pregunta1.pptx")
DOCX_PATH = os.path.join(BASE, "guion_presentacion.docx")
IMG_MATRIZ = os.path.join(BASE, "matriz_confusion.png")
IMG_IMP = os.path.join(BASE, "importancia_variables.png")
IMG_ROC = os.path.join(BASE, "curva_roc.png")

BLUE = RGBColor(0x1F, 0x4E, 0x8C)
GRAY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x22, 0x22, 0x22)

D_BLUE = DocxRGB(0x1F, 0x4E, 0x8C)
D_GRAY = DocxRGB(0x55, 0x55, 0x55)
D_BLACK = DocxRGB(0x22, 0x22, 0x22)


def add_textbox(slide, left, top, width, height, text, size=18, color=BLACK,
                bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, size=18, color=BLACK, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
        p.space_after = Pt(6)
    return tb


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 - Portada
    s1 = prs.slides.add_slide(blank)
    add_textbox(s1, Inches(0.6), Inches(0.5), Inches(12.0), Inches(1.0),
                "Predicción de bajo desempeño en inglés en Córdoba",
                size=36, color=BLUE, bold=True)
    add_textbox(s1, Inches(0.6), Inches(1.4), Inches(12.0), Inches(0.7),
                "Modelo predictivo para la Secretaría de Educación de Córdoba",
                size=20, color=GRAY)
    add_textbox(s1, Inches(0.6), Inches(2.5), Inches(12.0), Inches(2.0),
                "¿Podemos identificar — antes del Saber 11 — a los estudiantes que caerán en el nivel más bajo de inglés (A−), para focalizar refuerzo a tiempo?",
                size=22, color=BLUE, bold=True)
    add_bullets(s1, Inches(0.6), Inches(5.0), Inches(12.0), Inches(2.0), [
        "194,070 estudiantes de Córdoba analizados",
        "61% del departamento cae en nivel A− (sin alcanzar siquiera el básico inicial)",
        "Red neuronal entrenada con 137 variables del contexto familiar y del colegio",
    ], size=18, color=BLACK)

    # Slide 2 - Resultados
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9),
                "El modelo identifica al 95% de los estudiantes en riesgo",
                size=30, color=BLUE, bold=True)
    s2.shapes.add_picture(IMG_MATRIZ, Inches(0.5), Inches(1.5), height=Inches(5.4))
    # Right metrics
    add_textbox(s2, Inches(8.5), Inches(1.8), Inches(4.5), Inches(0.7),
                "Recall A−    95%", size=28, color=BLUE, bold=True)
    add_textbox(s2, Inches(8.5), Inches(2.7), Inches(4.5), Inches(0.7),
                "F1 score      0.79", size=28, color=BLUE, bold=True)
    add_textbox(s2, Inches(8.5), Inches(3.6), Inches(4.5), Inches(0.7),
                "AUC            0.74", size=28, color=BLUE, bold=True)
    add_textbox(s2, Inches(8.5), Inches(5.0), Inches(4.5), Inches(2.0),
                "De 17,751 estudiantes que cayeron en A−, el modelo identifica 16,831. Solo se le escapan 920.",
                size=14, color=GRAY)

    # Slide 3 - Hallazgo
    s3 = prs.slides.add_slide(blank)
    add_textbox(s3, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9),
                "El colegio pesa más que el estrato — el hallazgo clave",
                size=30, color=BLUE, bold=True)
    s3.shapes.add_picture(IMG_IMP, Inches(0.5), Inches(1.5), height=Inches(5.5))
    add_textbox(s3, Inches(8.5), Inches(1.8), Inches(4.5), Inches(1.5),
                "El colegio pesa 5x más que cualquier otra variable",
                size=22, color=BLUE, bold=True)
    add_bullets(s3, Inches(8.5), Inches(3.4), Inches(4.5), Inches(2.2), [
        "Más que el estrato",
        "Más que la educación de los padres",
        "Más que el municipio",
    ], size=16, color=BLACK)
    add_textbox(s3, Inches(8.5), Inches(5.5), Inches(4.5), Inches(1.8),
                "Internet y computador en casa pesan más que el estrato directo: el aprendizaje del inglés depende del consumo digital.",
                size=14, color=GRAY)

    # Slide 4 - Acciones
    s4 = prs.slides.add_slide(blank)
    add_textbox(s4, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9),
                "Tres frentes de acción para la Secretaría",
                size=30, color=BLUE, bold=True)

    blocks = [
        ("1. Identificar y replicar prácticas de colegios 'outperformers'",
         "Mapear qué hacen los colegios que sí logran sacar a sus estudiantes del A−. Es la palanca de mayor retorno."),
        ("2. Programas de conectividad digital sobre transferencias económicas",
         "Internet y computador en el hogar muestran más relación con el desempeño en inglés que las ayudas genéricas."),
        ("3. Estrategias diferenciadas por región",
         "Montería opera distinto al resto del departamento. No aplica una estrategia única para todo Córdoba."),
    ]
    top = 1.6
    for title, sub in blocks:
        add_textbox(s4, Inches(0.7), Inches(top), Inches(12.0), Inches(0.55),
                    title, size=20, color=BLUE, bold=True)
        add_textbox(s4, Inches(0.7), Inches(top + 0.55), Inches(12.0), Inches(0.7),
                    sub, size=15, color=GRAY)
        top += 1.5

    add_textbox(s4, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5),
                "Modelo serializado y listo para integrarse al tablero de la Secretaría.",
                size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)
    print("PPTX saved:", PPTX_PATH)


def add_doc_para(doc, text, size=11, bold=False, color=D_BLACK, font="Calibri",
                 align=None, space_after=6):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    run = p.add_run(text)
    run.font.size = DocxPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    p.paragraph_format.space_after = DocxPt(space_after)
    return p


def build_docx():
    doc = Document()
    # Title
    add_doc_para(doc,
                 "Guion de presentación — Pregunta 1: Bajo desempeño en inglés (A−)",
                 size=18, bold=True, color=D_BLUE)
    add_doc_para(doc,
                 "Duración objetivo: 2 minutos 30 segundos. Cliente: Secretaría de Educación de Córdoba.",
                 size=11, color=D_GRAY, space_after=12)

    slides = [
        ("Slide 1 — Portada y problema", "Tiempo: ~30 segundos",
         "Buenos días. Mi pregunta de negocio aborda una realidad preocupante de Córdoba: el 61% de los estudiantes del departamento sale del Saber 11 en el nivel más bajo de inglés, A−, sin alcanzar siquiera el básico inicial. La pregunta concreta que me planteé fue: ¿podemos identificar a esos estudiantes antes del examen, para que la Secretaría pueda intervenir a tiempo? Para responderla, entrené una red neuronal sobre 194,070 estudiantes del departamento, con 137 variables que combinan contexto familiar, características del colegio y del municipio."),
        ("Slide 2 — El modelo identifica al 95% de los estudiantes en riesgo", "Tiempo: ~40 segundos",
         "Esta es la prueba de fuego del modelo: lo evaluamos sobre 29,000 estudiantes que nunca había visto durante el entrenamiento. Y el resultado es contundente. De los 17,751 estudiantes que efectivamente cayeron en A−, el modelo identifica correctamente a 16,831 — el 95%. Solo se le escapan 920. Esto significa que prácticamente ningún estudiante en riesgo pasa desapercibido. El costo es que un porcentaje de los flagueados como riesgo termina alcanzando un nivel superior, pero para la Secretaría ese trade-off conviene: es mucho mejor incluir de más en programas de refuerzo que dejar fuera a un estudiante que sí lo necesitaba. El AUC de 0.74 confirma que el modelo es estadísticamente sólido y muy superior a una clasificación aleatoria."),
        ("Slide 3 — El hallazgo clave: el colegio pesa más que el estrato", "Tiempo: ~50 segundos",
         "Este es el hallazgo que más le interesa a la Secretaría. Cuando le preguntamos al modelo qué variables pesan más en su predicción, el resultado es contundente: el colegio al que asiste el estudiante explica cinco veces más que la siguiente variable. Más que el estrato, más que la educación de los padres, más que cualquier otro factor. Eso significa que en Córdoba existen colegios que sistemáticamente sacan a sus estudiantes del A−, y otros que sistemáticamente los dejan ahí, incluso controlando por el contexto socioeconómico. La heterogeneidad institucional es enorme. El segundo dato relevante: tener internet y computador en el hogar pesa más que el estrato directo. Esto tiene sentido: para aprender inglés moderno hace falta exposición al idioma, y esa exposición pasa hoy por la conectividad digital — series, música, herramientas de práctica."),
        ("Slide 4 — Tres frentes de acción y cierre", "Tiempo: ~30 segundos",
         "Esto se traduce en tres recomendaciones concretas para la Secretaría. Primero: identificar los colegios que sí están funcionando bien en inglés y mapear sus prácticas pedagógicas — esa es la palanca con mayor retorno. Segundo: priorizar programas de conectividad digital sobre transferencias económicas; los datos lo respaldan. Y tercero: las intervenciones deben ser regionales — lo que funciona en Montería no necesariamente funciona en el San Jorge o el Sinú medio. El modelo queda serializado y se integrará en el tablero que la Secretaría usará para focalizar intervenciones. Gracias."),
    ]

    for h1, h2, body in slides:
        add_doc_para(doc, h1, size=18, bold=True, color=D_BLUE, space_after=4)
        add_doc_para(doc, h2, size=13, bold=True, color=D_GRAY, space_after=4)
        add_doc_para(doc, body, size=11, color=D_BLACK, space_after=12)

    add_doc_para(doc, "Frases clave que NO debes olvidar", size=14, bold=True, color=D_BLUE, space_after=4)
    for line in [
        "95% de los estudiantes en riesgo identificados.",
        "El colegio pesa 5 veces más que cualquier otra variable.",
        "Internet y computador en casa pesan más que el estrato directo.",
    ]:
        p = doc.add_paragraph()
        run = p.add_run("• " + line)
        run.font.size = DocxPt(11)
        run.font.color.rgb = D_BLACK
        run.font.name = "Calibri"

    doc.save(DOCX_PATH)
    print("DOCX saved:", DOCX_PATH)


if __name__ == "__main__":
    build_pptx()
    build_docx()
    for p in (PPTX_PATH, DOCX_PATH):
        size = os.path.getsize(p)
        print(f"FILE {p}  SIZE {size} bytes")
