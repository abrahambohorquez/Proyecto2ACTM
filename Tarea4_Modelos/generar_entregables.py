# -*- coding: utf-8 -*-
"""Genera presentacion_pregunta1.pptx y guion_presentacion.docx."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

PPTX_PATH = os.path.join(BASE_DIR, "presentacion_pregunta1.pptx")
DOCX_PATH = os.path.join(BASE_DIR, "guion_presentacion.docx")

IMG_MATRIZ = os.path.join(BASE_DIR, "matriz_confusion.png")
IMG_IMPORT = os.path.join(BASE_DIR, "importancia_variables.png")
IMG_ROC = os.path.join(BASE_DIR, "curva_roc.png")

# Colores
AZUL = RGBColor(0x1F, 0x4E, 0x8C)
GRIS = RGBColor(0x55, 0x55, 0x55)
NEGRO = RGBColor(0x22, 0x22, 0x22)
AZUL_CLARO = RGBColor(0xD6, 0xE4, 0xF5)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

# Para docx
AZUL_DOCX = DocxRGB(0x1F, 0x4E, 0x8C)
GRIS_DOCX = DocxRGB(0x55, 0x55, 0x55)
NEGRO_DOCX = DocxRGB(0x22, 0x22, 0x22)


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False, anchor=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_bullets(slide, x, y, w, h, bullets, size, color, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(6)
    return tb


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---------- SLIDE 1 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.5, 12.1, 1.0,
                "Predicción de bajo desempeño en inglés en Córdoba",
                36, AZUL, bold=True)
    add_textbox(s, 0.6, 1.45, 12.1, 0.6,
                "Modelo predictivo para la Secretaría de Educación de Córdoba",
                20, GRIS)
    add_textbox(s, 0.6, 2.5, 12.1, 1.6,
                "¿Podemos identificar — antes del Saber 11 — a los estudiantes que caerán en el "
                "nivel más bajo de inglés (A−), para focalizar refuerzo a tiempo?",
                22, AZUL, bold=True)
    add_bullets(s, 0.8, 4.5, 11.8, 2.5, [
        "194,070 estudiantes de Córdoba analizados",
        "61% del departamento cae en nivel A− (sin alcanzar siquiera el básico inicial)",
        "137 variables del contexto familiar, del colegio y del municipio",
    ], 18, NEGRO)

    # ---------- SLIDE 2 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.4, 12.1, 0.9,
                "Metodología: rigor en el modelado", 30, AZUL, bold=True)

    bloques = [
        ("División de datos",
         "Split estratificado 70/15/15 — 135,926 train, 29,033 val, 29,111 test "
         "(test aislado hasta el final)"),
        ("Baseline obligatorio",
         "Regresión logística como referencia (F1 = 0.72, AUC = 0.70) — fija "
         "el piso que la red debe superar"),
        ("Red neuronal MLP — 5 arquitecturas probadas",
         "32-16  /  64-32  /  128-64  /  128-64-32  /  256-128-64 — con class weight "
         "balanceado y EarlyStopping"),
        ("Trazabilidad",
         "Todos los experimentos registrados en MLflow con sus parámetros, "
         "métricas y modelos serializados"),
    ]
    # 2x2 grid
    pos = [(0.6, 1.6), (6.95, 1.6), (0.6, 4.55), (6.95, 4.55)]
    for (titulo, desc), (x, y) in zip(bloques, pos):
        # caja con borde
        from pptx.enum.shapes import MSO_SHAPE
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(5.85), Inches(2.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFB)
        shape.line.color.rgb = AZUL
        shape.line.width = Pt(1.0)
        shape.text_frame.text = ""  # clear default
        # titulo
        add_textbox(s, x + 0.2, y + 0.15, 5.5, 0.6, titulo, 18, AZUL, bold=True)
        # descripcion
        add_textbox(s, x + 0.2, y + 0.85, 5.5, 1.7, desc, 14, NEGRO)

    # ---------- SLIDE 3 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.4, 12.1, 0.9,
                "Cómo construimos el mejor modelo: paso a paso",
                30, AZUL, bold=True)

    rows = 5
    cols = 3
    table_left = Inches(1.2)
    table_top = Inches(1.7)
    table_width = Inches(10.9)
    table_height = Inches(3.4)
    table = s.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
    table.columns[0].width = Inches(6.4)
    table.columns[1].width = Inches(2.25)
    table.columns[2].width = Inches(2.25)

    headers = ["Versión del modelo", "F1", "AUC"]
    data = [
        ("Baseline (regresión logística)", "0.72", "0.70", False),
        ("Red neuronal base (102 features)", "0.75", "0.71", False),
        ("+ Threshold optimizado (0.29)", "0.78", "0.71", False),
        ("+ Municipio + Colegio (target encoding)", "0.79", "0.74", True),
    ]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BLANCO
        run.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL

    for i, (v, f1, auc, highlight) in enumerate(data, start=1):
        for j, val in enumerate([v, f1, auc]):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(16)
            run.font.bold = highlight
            run.font.color.rgb = AZUL if highlight else NEGRO
            run.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = AZUL_CLARO if highlight else BLANCO

    add_textbox(s, 0.6, 5.6, 12.1, 1.2,
                "La mayor ganancia vino del feature engineering, no de cambiar la arquitectura",
                18, AZUL, bold=True, align=PP_ALIGN.CENTER)

    # ---------- SLIDE 4 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.4, 0.3, 12.5, 0.9,
                "El modelo identifica al 95% de los estudiantes en riesgo",
                30, AZUL, bold=True, align=PP_ALIGN.CENTER)
    # imagenes
    s.shapes.add_picture(IMG_MATRIZ, Inches(0.5), Inches(1.3), height=Inches(4.8))
    s.shapes.add_picture(IMG_ROC, Inches(7.0), Inches(1.3), height=Inches(4.8))

    # metricas
    add_textbox(s, 0.5, 6.25, 4.0, 0.7, "F1 = 0.79", 28, AZUL, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 4.7, 6.25, 4.0, 0.7, "Recall = 95%", 28, AZUL, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 8.9, 6.25, 4.0, 0.7, "AUC = 0.74", 28, AZUL, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(s, 0.4, 7.0, 12.5, 0.4,
                "Evaluación sobre 29,111 estudiantes que el modelo nunca había visto",
                12, GRIS, italic=True, align=PP_ALIGN.CENTER)

    # ---------- SLIDE 5 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.4, 0.3, 12.5, 0.9,
                "El colegio pesa más que el estrato — el hallazgo clave",
                30, AZUL, bold=True, align=PP_ALIGN.CENTER)
    s.shapes.add_picture(IMG_IMPORT, Inches(0.4), Inches(1.4), height=Inches(5.5))

    add_textbox(s, 7.6, 1.6, 5.4, 1.6,
                "El colegio pesa 5x más\nque cualquier otra variable",
                24, AZUL, bold=True)
    add_bullets(s, 7.7, 3.5, 5.3, 2.0, [
        "Más que el estrato",
        "Más que la educación de los padres",
        "Más que el municipio",
    ], 16, NEGRO)
    add_textbox(s, 7.6, 5.7, 5.4, 1.5,
                "Internet y computador en casa pesan más que el estrato directo: el aprendizaje "
                "del inglés depende del consumo digital.",
                14, GRIS, italic=True)

    # ---------- SLIDE 6 ----------
    s = prs.slides.add_slide(blank)
    add_textbox(s, 0.6, 0.3, 12.1, 0.8,
                "Acciones para la Secretaría", 28, AZUL, bold=True)

    recs = [
        ("1. Identificar y replicar prácticas de colegios 'outperformers'",
         "Mapear qué hacen los colegios que sí logran sacar a sus estudiantes del A−. "
         "Es la palanca de mayor retorno."),
        ("2. Conectividad digital sobre transferencias económicas",
         "Internet y computador en el hogar muestran más relación con el desempeño "
         "en inglés que las ayudas genéricas."),
        ("3. Estrategias regionales diferenciadas",
         "Montería opera distinto al resto del departamento. No aplica una estrategia "
         "única para todo Córdoba."),
    ]
    y = 1.2
    for titulo, desc in recs:
        add_textbox(s, 0.7, y, 12.0, 0.5, titulo, 18, AZUL, bold=True)
        add_textbox(s, 0.9, y + 0.5, 11.8, 0.7, desc, 14, GRIS)
        y += 1.25

    add_textbox(s, 0.6, y + 0.05, 12.1, 1.0,
                "Limitación: el AUC techo de 0.74 indica que hay variables fuera del dataset "
                "(motivación individual, calidad docente específica) que el modelo no captura. "
                "Se recomienda re-entrenar anualmente con datos nuevos.",
                14, GRIS, italic=True)

    add_textbox(s, 0.6, 7.05, 12.1, 0.4,
                "Modelo serializado y listo para integrarse al tablero de la Secretaría.",
                14, AZUL, bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)
    print("PPTX guardado en", PPTX_PATH)


# -----------------------------------------------------------------------------
def set_run_format(run, size, color, bold=False, italic=False, name="Calibri"):
    run.font.name = name
    run.font.size = DocxPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_para(doc, text, size, color, bold=False, italic=False, align=None):
    p = doc.add_paragraph()
    if align is not None:
        p.alignment = align
    run = p.add_run(text)
    set_run_format(run, size, color, bold=bold, italic=italic)
    return p


def build_docx():
    doc = Document()

    # Fuente por defecto
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)
    style.font.color.rgb = NEGRO_DOCX

    add_para(doc,
             "Guion de presentación — Pregunta 1: Bajo desempeño en inglés (A−)",
             18, AZUL_DOCX, bold=True)
    add_para(doc,
             "Duración objetivo: 2 minutos 40 segundos. Cliente: Secretaría de Educación de Córdoba.",
             12, GRIS_DOCX)
    doc.add_paragraph("")

    secciones = [
        ("Slide 1 — Portada y problema", "Tiempo: ~25 segundos",
         "Buenos días. El 61% de los estudiantes de Córdoba sale del Saber 11 en el nivel "
         "más bajo de inglés, A−, sin alcanzar siquiera el básico inicial. Mi pregunta "
         "de negocio fue: ¿podemos identificar a esos estudiantes antes del examen, para que la "
         "Secretaría pueda intervenir a tiempo? Para responderla, trabajé sobre los 194,070 "
         "estudiantes del departamento con 137 variables que combinan contexto familiar, "
         "características del colegio y municipio."),
        ("Slide 2 — Metodología", "Tiempo: ~25 segundos",
         "Para construir el modelo seguí tres pilares de rigor. Primero, dividí los datos en "
         "tres conjuntos estratificados: 70% para entrenar, 15% para validación y 15% para test, "
         "este último intacto durante todo el desarrollo. Segundo, antes de meter una red neuronal "
         "entrené una regresión logística como baseline; alcanzó F1 de 0.72. Eso "
         "fijó el piso que la red tenía que superar. Tercero, probé cinco arquitecturas "
         "distintas de MLP — desde una red pequeña 32-16 hasta una grande 256-128-64 — "
         "todas con class weight balanceado y EarlyStopping. Cada experimento quedó registrado "
         "en MLflow para trazabilidad y reproducibilidad."),
        ("Slide 3 — Journey de mejora", "Tiempo: ~25 segundos",
         "Esta es la evolución completa del modelo. Arranqué en F1 de 0.72 con la "
         "regresión logística. La red neuronal con 102 features la superó hasta 0.75. "
         "Optimizando el threshold de decisión a 0.29, alineado con el costo asimétrico del "
         "cliente, subí a 0.78. Pero el salto definitivo vino al hacer feature engineering: "
         "agregué el municipio y el colegio del estudiante mediante target encoding con "
         "suavizado bayesiano. Eso disparó el AUC de 0.71 a 0.74 y consolidó el F1 en 0.79. "
         "La conclusión metodológica es clara: la mayor ganancia no vino de cambiar la "
         "arquitectura, sino de darle al modelo features con más señal."),
        ("Slide 4 — Resultados finales", "Tiempo: ~30 segundos",
         "Estos son los resultados sobre los 29,000 estudiantes de test que el modelo nunca había "
         "visto. La matriz de confusión a la izquierda muestra que de los 17,751 estudiantes que "
         "efectivamente cayeron en A−, el modelo identifica correctamente a 16,831 — el 95%. "
         "Solo se le escapan 920. La curva ROC a la derecha confirma que el modelo es "
         "estadísticamente sólido: un AUC de 0.74 lo coloca muy por encima de un clasificador "
         "aleatorio. El costo de la alta sensibilidad es un porcentaje de falsos positivos, pero ese "
         "trade-off conviene a la Secretaría: es preferible incluir de más en programas de "
         "refuerzo que dejar fuera a quien lo necesita."),
        ("Slide 5 — Hallazgo clave", "Tiempo: ~35 segundos",
         "Este es el hallazgo que más le interesa a la Secretaría. Cuando le preguntamos al "
         "modelo qué variables pesan más en su predicción, el resultado es contundente: "
         "el colegio al que asiste el estudiante explica cinco veces más que la siguiente "
         "variable. Más que el estrato, más que la educación de los padres, más "
         "que cualquier otro factor. Eso significa que en Córdoba existen colegios que "
         "sistemáticamente sacan a sus estudiantes del A−, y otros que sistemáticamente "
         "los dejan ahí, incluso controlando por el contexto socioeconómico. Es una "
         "heterogeneidad institucional enorme. El segundo dato relevante: tener internet y computador "
         "en el hogar pesa más que el estrato directo. Esto tiene sentido porque aprender "
         "inglés moderno depende del consumo digital."),
        ("Slide 6 — Recomendaciones, limitaciones y cierre", "Tiempo: ~20 segundos",
         "Esto se traduce en tres recomendaciones para la Secretaría. Primero: identificar y "
         "replicar las prácticas pedagógicas de los colegios que están funcionando. "
         "Segundo: priorizar conectividad digital sobre transferencias económicas. Tercero: "
         "estrategias regionales diferenciadas — Montería opera distinto al resto del "
         "departamento. Como limitación honesta, el AUC techo de 0.74 indica que hay variables "
         "fuera del dataset, como motivación individual o calidad docente específica, que el "
         "modelo no captura. El modelo queda serializado y listo para integrarse al tablero. Gracias."),
    ]

    for h1, h2, parrafo in secciones:
        add_para(doc, h1, 16, AZUL_DOCX, bold=True)
        add_para(doc, h2, 12, GRIS_DOCX, bold=True)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        run = p.add_run(parrafo)
        set_run_format(run, 11, NEGRO_DOCX)
        doc.add_paragraph("")

    add_para(doc, "Frases clave que NO debes olvidar", 14, AZUL_DOCX, bold=True)
    bullets = [
        "95% de los estudiantes en riesgo identificados.",
        "El colegio pesa 5 veces más que cualquier otra variable.",
        "Internet y computador en casa pesan más que el estrato directo.",
        "La mayor ganancia vino del feature engineering, no de cambiar la arquitectura.",
    ]
    for b in bullets:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(b)
        set_run_format(run, 11, NEGRO_DOCX)

    doc.save(DOCX_PATH)
    print("DOCX guardado en", DOCX_PATH)


if __name__ == "__main__":
    for img in [IMG_MATRIZ, IMG_IMPORT, IMG_ROC]:
        if not os.path.exists(img):
            raise SystemExit("FALTA imagen: " + img)
    build_pptx()
    build_docx()
    for p in [PPTX_PATH, DOCX_PATH]:
        size = os.path.getsize(p)
        print("OK:", p, size, "bytes")
